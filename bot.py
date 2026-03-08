"""
🏫 ПроУрок v1.3 — Telegram-бот помощник для учителей
Исправления: Markdown баги, баланс по кнопке, убрана кнопка Материал
Новое: улучшенные презентации, диаграммы Excel, интерактивные HTML тесты
"""

import os
import io
import json
import logging
import tempfile
import re
import base64
import subprocess
import speech_recognition as sr
from pydub import AudioSegment
from datetime import datetime
from pathlib import Path
import PyPDF2
from docx import Document as DocxDocument

import requests
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np

from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup, ReplyKeyboardMarkup, KeyboardButton
from telegram.ext import (
    Application, CommandHandler, MessageHandler,
    CallbackQueryHandler, ContextTypes, filters
)
from anthropic import Anthropic

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, PieChart, LineChart, Reference

from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

logging.basicConfig(format="%(asctime)s - %(name)s - %(levelname)s - %(message)s", level=logging.INFO)
logger = logging.getLogger(__name__)

TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")
client = Anthropic(api_key=ANTHROPIC_API_KEY)

BOT_NAME = "ПроУрок"
BOT_VERSION = "1.3"
SUPPORT_BOT_URL = "https://t.me/pro_lessen_help_bot"

TARIFFS = {
    "demo": {"name": "Демо", "queries": 20, "generations": 5, "price": "Бесплатно"},
    "start": {"name": "Старт", "queries": 200, "generations": 30, "price": "490 руб/мес"},
    "pro": {"name": "Про", "queries": 500, "generations": 100, "price": "890 руб/мес"},
    "premium": {"name": "Премиум", "queries": 999999, "generations": 999999, "price": "1490 руб/мес"},
}

LIMITS_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "user_limits.json")

def load_all_limits():
    if os.path.exists(LIMITS_FILE):
        try:
            with open(LIMITS_FILE, 'r', encoding='utf-8') as f: return json.load(f)
        except: return {}
    return {}

def save_all_limits(data):
    try:
        with open(LIMITS_FILE, 'w', encoding='utf-8') as f: json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e: logger.error(f"Save: {e}")

def get_user_data(user_id):
    d = load_all_limits(); uid = str(user_id)
    if uid not in d:
        d[uid] = {"tariff": "demo", "queries_used": 0, "generations_used": 0,
                  "registered": datetime.now().strftime("%d.%m.%Y"), "name": "", "username": ""}
        save_all_limits(d)
    return d[uid]

def update_user_data(user_id, data):
    d = load_all_limits(); d[str(user_id)] = data; save_all_limits(d)

def save_user_info(user_id, first_name, username):
    data = get_user_data(user_id)
    data["name"] = first_name or ""; data["username"] = username or ""
    update_user_data(user_id, data)

def get_tariff(user_id): return TARIFFS.get(get_user_data(user_id)["tariff"], TARIFFS["demo"])

def check_query_limit(user_id):
    d = get_user_data(user_id); t = TARIFFS.get(d["tariff"], TARIFFS["demo"])
    return d["queries_used"] < t["queries"]

def check_generation_limit(user_id):
    d = get_user_data(user_id); t = TARIFFS.get(d["tariff"], TARIFFS["demo"])
    return d["generations_used"] < t["generations"]

def use_query(user_id):
    d = get_user_data(user_id); d["queries_used"] += 1; update_user_data(user_id, d)

def use_generation(user_id):
    d = get_user_data(user_id); d["generations_used"] += 1; update_user_data(user_id, d)

def get_limit_status(user_id):
    d = get_user_data(user_id); t = TARIFFS.get(d["tariff"], TARIFFS["demo"])
    return max(0, t["queries"] - d["queries_used"]), max(0, t["generations"] - d["generations_used"])

def get_balance_text(user_id):
    d = get_user_data(user_id); t = TARIFFS.get(d["tariff"], TARIFFS["demo"])
    ql, gl = get_limit_status(user_id)
    text = "Ваш баланс:\n\n"
    text += f"Тариф: {t['name']} ({t['price']})\n\n"
    if d["tariff"] == "premium":
        text += "Запросов: безлимит\nОформлений: безлимит\n"
    else:
        text += f"Запросов: {d['queries_used']} / {t['queries']} (осталось {ql})\n"
        text += f"Оформлений: {d['generations_used']} / {t['generations']} (осталось {gl})\n"
    if d["tariff"] == "demo" and ql < 5:
        text += f"\nЗапросы заканчиваются!\nПодключите тариф: {SUPPORT_BOT_URL}"
    elif d["tariff"] != "premium" and ql < 20:
        text += f"\nЗапросы заканчиваются. Повысьте тариф: {SUPPORT_BOT_URL}"
    return text

LIMIT_MESSAGE = (
    "Лимит исчерпан!\n\n"
    f"Подключите тариф:\n{SUPPORT_BOT_URL}\n\n"
    "Старт — 490 руб/мес (200 запросов)\n"
    "Про — 890 руб/мес (500 запросов)\n"
    "Премиум — 1490 руб/мес (безлимит)"
)

GENERATION_LIMIT_MESSAGE = (
    "Лимит оформления документов исчерпан!\n\n"
    "Текстовые запросы ещё доступны.\n\n"
    f"Повысьте тариф:\n{SUPPORT_BOT_URL}"
)

user_conversations = {}
MAX_HISTORY = 20

SYSTEM_PROMPT = f"""Ты — профессиональный помощник учителя. Твоё имя — {BOT_NAME}.
Ты помогаешь учителям ЛЮБЫХ предметов: математика, русский язык, литература, история, обществознание, физика, химия, биология, география, английский язык, информатика, ИЗО, музыка, физкультура, технология, ОБЖ и другие.

ВОЗМОЖНОСТИ: планы уроков (1-11 класс), тесты, конспекты, рабочие листы, ФГОС, ОГЭ/ЕГЭ, анализ фото, идеи для уроков.

ПРАВИЛА:
- Русский язык, развёрнуто, с примерами
- Адаптируй под класс и предмет
- НЕ используй маркеры [PPTX_DOC], [WORD_DOC], [EXCEL_DOC]
- Просто давай качественный текстовый ответ

Источники: resh.edu.ru, yaklass.ru, infourok.ru, fipi.ru
"""

PPTX_PROMPT = """
Переработай материал в презентацию. Начни с [PPTX_DOC].
Формат: [SLIDE]ЗАГОЛОВОК:..ПОДЗАГОЛОВОК:..СОДЕРЖИМОЕ:-пункт..ЗАМЕТКИ:..[/SLIDE]
Диаграммы: [CHART]ТИП:bar/pie/line НАЗВАНИЕ:.. МЕТКИ:.. ЗНАЧЕНИЯ:.. ЕДИНИЦЫ:..[/CHART]
Картинки: [IMAGE]ЗАПРОС:english query ОПИСАНИЕ:..[/IMAGE]
10 слайдов, 5-7 развёрнутых пунктов, минимум 10-15 слов каждый.
"""

WORD_PROMPT = "Переработай в документ. Начни с [WORD_DOC]. Заголовки (#,##,###), списки, структура."
EXCEL_PROMPT = """Переработай в таблицу. Начни с [EXCEL_DOC]. Столбцы через |.
Если данные числовые — добавь маркер для диаграммы:
[EXCEL_CHART]
ТИП: bar (или pie, line)
НАЗВАНИЕ: Название диаграммы
[/EXCEL_CHART]
"""

HTML_TEST_PROMPT = """Создай интерактивный HTML тест для учеников.
Начни ответ с маркера [HTML_TEST].
Сделай полноценный HTML файл с встроенными CSS и JavaScript:

Требования к тесту:
- 10-15 вопросов с вариантами ответа (4 варианта, 1 правильный)
- Красивое современное оформление (градиенты, тени, анимации)
- Яркий дизайн привлекательный для учеников
- При выборе неправильного ответа — подсветить красным и показать правильный
- При правильном ответе — подсветить зелёным с анимацией
- Кнопка «Следующий вопрос» после ответа
- В конце — результат с оценкой (5/4/3/2) и процентом
- Возможность пройти заново
- Адаптивный дизайн (работает на телефоне)
- Всё на русском языке
- Шрифт: Inter или системный sans-serif
- Цвета: современная палитра с градиентами

Верни ТОЛЬКО HTML код, без пояснений, без markdown блоков.
"""

def get_history(uid):
    if uid not in user_conversations: user_conversations[uid] = []
    return user_conversations[uid]

def add_to_history(uid, role, content):
    h = get_history(uid); h.append({"role": role, "content": content})
    if len(h) > MAX_HISTORY: user_conversations[uid] = h[-MAX_HISTORY:]

def ask_claude(uid, msg, image_data=None):
    if image_data:
        content = [{"type": "image", "source": {"type": "base64", "media_type": image_data['mime'], "data": image_data['base64']}}, {"type": "text", "text": msg}]
        add_to_history(uid, "user", msg); messages = [{"role": "user", "content": content}]
    else:
        add_to_history(uid, "user", msg); messages = get_history(uid)
    try:
        r = client.messages.create(model="claude-3-haiku-20240307", max_tokens=4096, system=SYSTEM_PROMPT, messages=messages)
        answer = r.content[0].text; add_to_history(uid, "assistant", answer); return answer
    except Exception as e:
        logger.error(f"Claude: {e}"); return f"Ошибка ИИ: {e}"

def clean_for_chat(text):
    c = text
    for m in ['[PPTX_DOC]','[WORD_DOC]','[EXCEL_DOC]','[HTML_TEST]']: c = c.replace(m, '')
    c = re.sub(r'\[/?SLIDE\]', '', c); c = re.sub(r'\[CHART\].*?\[/CHART\]', '', c, flags=re.DOTALL)
    c = re.sub(r'\[IMAGE\].*?\[/IMAGE\]', '', c, flags=re.DOTALL)
    c = re.sub(r'\[EXCEL_CHART\].*?\[/EXCEL_CHART\]', '', c, flags=re.DOTALL)
    return re.sub(r'\n{3,}', '\n\n', c).strip()

async def safe_reply(update, text):
    """Отправка без Markdown — безопасно от спецсимволов."""
    if len(text) > 4000:
        for i in range(0, len(text), 4000):
            await update.message.reply_text(text[i:i+4000])
    else:
        await update.message.reply_text(text)


# ─── Картинки ────────────────────────────────────────────────

def download_image(query):
    try:
        url = "https://commons.wikimedia.org/w/api.php"
        r = requests.get(url, params={"action":"query","list":"search","srsearch":query,"srnamespace":"6","srlimit":"3","format":"json"}, timeout=10)
        d = r.json()
        if "query" in d and d["query"]["search"]:
            t = d["query"]["search"][0]["title"]
            r2 = requests.get(url, params={"action":"query","titles":t,"prop":"imageinfo","iiprop":"url|mime","format":"json"}, timeout=10)
            d2 = r2.json()
            if "query" in d2 and "pages" in d2["query"]:
                for pid, pd in d2["query"]["pages"].items():
                    if "imageinfo" in pd and "image" in pd["imageinfo"][0].get("mime",""):
                        ir = requests.get(pd["imageinfo"][0]["url"], timeout=15)
                        if ir.status_code == 200: return io.BytesIO(ir.content)
    except: pass
    try:
        r = requests.get(f"https://source.unsplash.com/800x500/?{query}", timeout=15, allow_redirects=True)
        if r.status_code == 200 and len(r.content) > 1000: return io.BytesIO(r.content)
    except: pass
    return None


# ─── Диаграммы matplotlib ─────────────────────────────────────

def create_chart_image(ct, title, labels, values, units="", tc=None):
    if not tc: tc = ['#4A90D9','#50C878','#FF6B6B','#FFD93D','#6C5CE7','#A8E6CF','#FF8A80','#82B1FF']
    plt.rcParams['font.size']=14; fig,ax=plt.subplots(figsize=(8,5))
    if ct=='bar':
        bars=ax.bar(labels,values,color=tc[:len(labels)],edgecolor='white')
        for b,v in zip(bars,values): ax.text(b.get_x()+b.get_width()/2.,b.get_height()+max(values)*0.02,f'{v:g}',ha='center',va='bottom',fontsize=12,fontweight='bold')
        ax.set_ylabel(units); ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        if any(len(str(l))>8 for l in labels): plt.xticks(rotation=30,ha='right')
    elif ct=='pie': ax.pie(values,labels=labels,autopct='%1.1f%%',colors=tc[:len(labels)],startangle=90,wedgeprops={'edgecolor':'white','linewidth':2})
    elif ct=='line':
        ax.plot(labels,values,color=tc[0],linewidth=3,marker='o',markersize=8,markerfacecolor=tc[1])
        ax.fill_between(range(len(labels)),values,alpha=0.1,color=tc[0]); ax.set_ylabel(units); ax.grid(axis='y',alpha=0.3)
        if any(len(str(l))>8 for l in labels): plt.xticks(rotation=30,ha='right')
    ax.set_title(title,fontsize=16,fontweight='bold',pad=15,color='#1E293B'); plt.tight_layout()
    buf=io.BytesIO(); fig.savefig(buf,format='png',dpi=150,bbox_inches='tight',facecolor='white'); plt.close(fig); buf.seek(0); return buf


# ─── Темы ─────────────────────────────────────────────────────

THEMES = {
    "blue":{"bg_dark":PptxRGB(26,54,93),"bg_light":PptxRGB(240,244,250),"accent":PptxRGB(74,144,217),"text_light":PptxRGB(255,255,255),"text_dark":PptxRGB(30,41,59),"subtitle":PptxRGB(180,200,230),"chart_colors":['#4A90D9','#50C878','#FF6B6B','#FFD93D','#6C5CE7']},
    "green":{"bg_dark":PptxRGB(34,87,60),"bg_light":PptxRGB(242,248,244),"accent":PptxRGB(80,200,120),"text_light":PptxRGB(255,255,255),"text_dark":PptxRGB(30,45,35),"subtitle":PptxRGB(180,220,195),"chart_colors":['#22573C','#50C878','#A8E6CF','#4A90D9','#FFD93D']},
    "warm":{"bg_dark":PptxRGB(93,52,40),"bg_light":PptxRGB(250,245,240),"accent":PptxRGB(217,144,74),"text_light":PptxRGB(255,255,255),"text_dark":PptxRGB(60,40,35),"subtitle":PptxRGB(230,200,180),"chart_colors":['#D9904A','#C86450','#FFD93D','#A8E6CF','#6C5CE7']},
}
def choose_theme(t):
    tl=t.lower()
    if any(w in tl for w in ['матем','физик','информ','алгебр']): return THEMES["blue"]
    elif any(w in tl for w in ['биолог','природ','географ']): return THEMES["green"]
    elif any(w in tl for w in ['истор','литератур','музык']): return THEMES["warm"]
    return THEMES["blue"]


# ─── Парсинг ─────────────────────────────────────────────────

def parse_chart(text):
    charts=[]
    for m in re.findall(r'\[CHART\](.*?)\[/CHART\]',text,re.DOTALL):
        c={'type':'bar','title':'','labels':[],'values':[],'units':''}
        for l in m.strip().split('\n'):
            l=l.strip()
            if l.startswith('ТИП:'): c['type']=l[4:].strip().lower()
            elif l.startswith('НАЗВАНИЕ:'): c['title']=l[9:].strip()
            elif l.startswith('МЕТКИ:'): c['labels']=[x.strip() for x in l[6:].split(',')]
            elif l.startswith('ЗНАЧЕНИЯ:'):
                try: c['values']=[float(x.strip().replace(' ','')) for x in l[9:].split(',')]
                except: pass
            elif l.startswith('ЕДИНИЦЫ:'): c['units']=l[8:].strip()
        if c['labels'] and c['values']: charts.append(c)
    return charts

def parse_image(text):
    imgs=[]
    for m in re.findall(r'\[IMAGE\](.*?)\[/IMAGE\]',text,re.DOTALL):
        img={'query':'','description':''}
        for l in m.strip().split('\n'):
            l=l.strip()
            if l.startswith('ЗАПРОС:'): img['query']=l[7:].strip()
            elif l.startswith('ОПИСАНИЕ:'): img['description']=l[9:].strip()
        if img['query']: imgs.append(img)
    return imgs

def parse_excel_chart(text):
    for m in re.findall(r'\[EXCEL_CHART\](.*?)\[/EXCEL_CHART\]',text,re.DOTALL):
        c={'type':'bar','title':'Диаграмма'}
        for l in m.strip().split('\n'):
            l=l.strip()
            if l.startswith('ТИП:'): c['type']=l[4:].strip().lower()
            elif l.startswith('НАЗВАНИЕ:'): c['title']=l[9:].strip()
        return c
    return None

def parse_slides(content):
    slides=[]
    for part in content.split('[SLIDE]'):
        if '[/SLIDE]' not in part: continue
        st=part.split('[/SLIDE]')[0].strip()
        sl={'title':'','subtitle':'','content':[],'notes':'','charts':parse_chart(st),'images':parse_image(st)}
        clean=re.sub(r'\[CHART\].*?\[/CHART\]','',st,flags=re.DOTALL)
        clean=re.sub(r'\[IMAGE\].*?\[/IMAGE\]','',clean,flags=re.DOTALL)
        sec=None
        for l in clean.split('\n'):
            s=l.strip()
            if s.startswith('ЗАГОЛОВОК:'): sl['title']=s[10:].strip()
            elif s.startswith('ПОДЗАГОЛОВОК:'): sl['subtitle']=s[13:].strip()
            elif s.startswith('СОДЕРЖИМОЕ:'): sec='content'
            elif s.startswith('ЗАМЕТКИ:'): sl['notes']=s[8:].strip(); sec='notes'
            elif sec=='content' and s: sl['content'].append(s[2:] if s[:2] in ['- ','• '] else s)
            elif sec=='notes' and s: sl['notes']+=' '+s
        if sl['title']: slides.append(sl)
    return slides

def create_slides_from_text(title,content):
    slides=[{'title':title,'subtitle':'','content':[],'notes':'','charts':[],'images':[]}]
    clean=re.sub(r'\[CHART\].*?\[/CHART\]','',content,flags=re.DOTALL)
    clean=re.sub(r'\[IMAGE\].*?\[/IMAGE\]','',clean,flags=re.DOTALL)
    cur=None
    for l in clean.split('\n'):
        s=l.strip()
        if not s: continue
        if s.startswith('# ') or s.startswith('## '): cur={'title':s.lstrip('#').strip(),'subtitle':'','content':[],'notes':'','charts':[],'images':[]}; slides.append(cur)
        elif cur:
            if s[:2] in ['- ','• ']: cur['content'].append(s[2:])
            elif len(s)>5: cur['content'].append(s)
    return slides if len(slides)>1 else [{'title':title,'subtitle':'','content':[content[:200]],'notes':'','charts':[],'images':[]}]


# ─── PowerPoint ───────────────────────────────────────────────

def add_slide_number(sl,num,theme,sw,sh):
    bg=sl.background;fill=bg.fill;fill.solid();fill.fore_color.rgb=theme['bg_light']
    bar=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,PptxInches(0.3),sh);bar.fill.solid();bar.fill.fore_color.rgb=theme['bg_dark'];bar.line.fill.background()
    acc=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,PptxInches(0.3),0,PptxInches(0.05),sh);acc.fill.solid();acc.fill.fore_color.rgb=theme['accent'];acc.line.fill.background()
    ns=sl.shapes.add_shape(MSO_SHAPE.OVAL,PptxInches(0.65),PptxInches(0.35),PptxInches(0.55),PptxInches(0.55));ns.fill.solid();ns.fill.fore_color.rgb=theme['bg_dark'];ns.line.fill.background()
    ntf=ns.text_frame;ntf.paragraphs[0].text=str(num);ntf.paragraphs[0].font.size=PptxPt(18);ntf.paragraphs[0].font.bold=True;ntf.paragraphs[0].font.color.rgb=theme['text_light'];ntf.paragraphs[0].alignment=PP_ALIGN.CENTER

def create_pptx_presentation(title, content, user_photo=None):
    sd=parse_slides(content)
    if not sd: sd=create_slides_from_text(title,content)
    theme=choose_theme(title);prs=Presentation();prs.slide_width=PptxInches(13.333);prs.slide_height=PptxInches(7.5);sw=prs.slide_width;sh=prs.slide_height
    # ТИТУЛЬНЫЙ
    sl=prs.slides.add_slide(prs.slide_layouts[6]);bg=sl.background;fill=bg.fill;fill.solid();fill.fore_color.rgb=theme['bg_dark']
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,sw,PptxInches(0.12));s.fill.solid();s.fill.fore_color.rgb=theme['accent'];s.line.fill.background()
    d=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,PptxInches(0.8),PptxInches(1.5),PptxInches(0.12),PptxInches(3));d.fill.solid();d.fill.fore_color.rgb=theme['accent'];d.line.fill.background()
    ft=sd[0]['title'] if sd else title
    tb=sl.shapes.add_textbox(PptxInches(1.3),PptxInches(1.8),PptxInches(10),PptxInches(2));tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text=ft;p.font.size=PptxPt(44);p.font.bold=True;p.font.color.rgb=theme['text_light'];p.font.name="Georgia"
    if sd and sd[0].get('subtitle'):
        p2=tf.add_paragraph();p2.text=sd[0]['subtitle'];p2.font.size=PptxPt(22);p2.font.color.rgb=theme['subtitle'];p2.font.name="Calibri"
    tb2=sl.shapes.add_textbox(PptxInches(1.3),PptxInches(5.5),PptxInches(8),PptxInches(0.8))
    p3=tb2.text_frame.paragraphs[0];p3.text=f"Подготовлено в {BOT_NAME}  {datetime.now().strftime('%d.%m.%Y')}";p3.font.size=PptxPt(14);p3.font.color.rgb=theme['subtitle']
    # Фото пользователя на титульный
    if user_photo:
        try: sl.shapes.add_picture(user_photo, PptxInches(8.5), PptxInches(1.5), PptxInches(4), PptxInches(4))
        except: pass
    elif sd and sd[0].get('images'):
        img=download_image(sd[0]['images'][0]['query'])
        if img:
            try: sl.shapes.add_picture(img,PptxInches(8.5),PptxInches(1.5),PptxInches(4),PptxInches(4))
            except: pass
    s2=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,sh-PptxInches(0.12),sw,PptxInches(0.12));s2.fill.solid();s2.fill.fore_color.rgb=theme['accent'];s2.line.fill.background()
    # СОДЕРЖИМОЕ
    cs=sd[1:] if len(sd)>1 else sd
    for i,sdata in enumerate(cs):
        sl=prs.slides.add_slide(prs.slide_layouts[6]);add_slide_number(sl,i+1,theme,sw,sh)
        hv=bool(sdata.get('charts') or sdata.get('images'));cw=PptxInches(6) if hv else PptxInches(11)
        tb=sl.shapes.add_textbox(PptxInches(1.5),PptxInches(0.3),PptxInches(11),PptxInches(0.8));ttf=tb.text_frame;ttf.word_wrap=True
        tp=ttf.paragraphs[0];tp.text=sdata['title'];tp.font.size=PptxPt(30);tp.font.bold=True;tp.font.color.rgb=theme['text_dark'];tp.font.name="Georgia"
        ln=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,PptxInches(1.5),PptxInches(1.15),PptxInches(2.5),PptxInches(0.04));ln.fill.solid();ln.fill.fore_color.rgb=theme['accent'];ln.line.fill.background()
        if sdata['content']:
            cb=sl.shapes.add_textbox(PptxInches(1.5),PptxInches(1.5),cw,PptxInches(5));ctf=cb.text_frame;ctf.word_wrap=True
            for j,item in enumerate(sdata['content']):
                pp=ctf.paragraphs[0] if j==0 else ctf.add_paragraph();pp.text=f"  {item}";pp.font.size=PptxPt(16);pp.font.color.rgb=theme['text_dark'];pp.font.name="Calibri";pp.space_after=PptxPt(10)
        if sdata.get('charts'):
            cd=sdata['charts'][0];ci=create_chart_image(cd['type'],cd['title'],cd['labels'],cd['values'],cd.get('units',''),theme.get('chart_colors'))
            if ci: sl.shapes.add_picture(ci,PptxInches(7.5),PptxInches(1.3),PptxInches(5.3),PptxInches(3.5))
        elif sdata.get('images'):
            img=download_image(sdata['images'][0]['query'])
            if img:
                try:
                    sl.shapes.add_picture(img,PptxInches(7.8),PptxInches(1.3),PptxInches(4.8),PptxInches(3.5))
                    cap=sl.shapes.add_textbox(PptxInches(7.8),PptxInches(4.9),PptxInches(4.8),PptxInches(0.5))
                    cp=cap.text_frame.paragraphs[0];cp.text=sdata['images'][0].get('description','');cp.font.size=PptxPt(10);cp.font.italic=True;cp.font.color.rgb=PptxRGB(140,140,140);cp.alignment=PP_ALIGN.CENTER
                except: pass
        if sdata.get('notes'): sl.notes_slide.notes_text_frame.text=sdata['notes']
        ft2=sl.shapes.add_textbox(PptxInches(0.7),PptxInches(6.95),PptxInches(12),PptxInches(0.35))
        fp=ft2.text_frame.paragraphs[0];fp.text=f"{BOT_NAME}  {sdata['title'][:50]}";fp.font.size=PptxPt(9);fp.font.color.rgb=PptxRGB(170,170,170)
    # ФИНАЛЬНЫЙ
    sl=prs.slides.add_slide(prs.slide_layouts[6]);bg=sl.background;fill=bg.fill;fill.solid();fill.fore_color.rgb=theme['bg_dark']
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,sw,PptxInches(0.12));s.fill.solid();s.fill.fore_color.rgb=theme['accent'];s.line.fill.background()
    tb=sl.shapes.add_textbox(PptxInches(1),PptxInches(2.2),PptxInches(11.333),PptxInches(2.5));tf=tb.text_frame;tf.word_wrap=True
    p=tf.paragraphs[0];p.text="Спасибо за внимание!";p.font.size=PptxPt(44);p.font.bold=True;p.font.color.rgb=theme['text_light'];p.font.name="Georgia";p.alignment=PP_ALIGN.CENTER
    p2=tf.add_paragraph();p2.text=f"Создано в {BOT_NAME}";p2.font.size=PptxPt(18);p2.font.color.rgb=theme['subtitle'];p2.alignment=PP_ALIGN.CENTER
    s2=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,sh-PptxInches(0.12),sw,PptxInches(0.12));s2.fill.solid();s2.fill.fore_color.rgb=theme['accent'];s2.line.fill.background()
    buf=io.BytesIO();prs.save(buf);buf.seek(0);return buf


# ─── Word ─────────────────────────────────────────────────────

def create_word_document(title,content):
    doc=Document();doc.styles['Normal'].font.name='Times New Roman';doc.styles['Normal'].font.size=Pt(12)
    h=doc.add_heading(title,level=0);h.alignment=WD_ALIGN_PARAGRAPH.CENTER
    for r in h.runs: r.font.color.rgb=RGBColor(26,54,93)
    dp=doc.add_paragraph();dp.alignment=WD_ALIGN_PARAGRAPH.RIGHT;dr=dp.add_run(f"Дата: {datetime.now().strftime('%d.%m.%Y')}");dr.font.size=Pt(10);dr.font.color.rgb=RGBColor(128,128,128)
    doc.add_paragraph()
    for line in content.split('\n'):
        s=line.strip()
        if not s: doc.add_paragraph()
        elif s.startswith('## '): doc.add_heading(s[3:],level=2)
        elif s.startswith('### '): doc.add_heading(s[4:],level=3)
        elif s.startswith('# '): doc.add_heading(s[2:],level=1)
        elif s[:2] in ['- ','• ']: doc.add_paragraph(s[2:],style='List Bullet')
        else:
            p=doc.add_paragraph();parts=s.split('**')
            for idx,part in enumerate(parts):
                if part: run=p.add_run(part);run.bold=(idx%2==1)
    buf=io.BytesIO();doc.save(buf);buf.seek(0);return buf


# ─── Excel с диаграммами ──────────────────────────────────────

def create_excel_document(title, content):
    wb = openpyxl.Workbook(); ws = wb.active; ws.title = title[:31]
    hf = Font(name='Arial',size=12,bold=True,color='FFFFFF')
    hfl = PatternFill(start_color='1A365D',end_color='1A365D',fill_type='solid')
    ha = Alignment(horizontal='center',vertical='center',wrap_text=True)
    cf = Font(name='Arial',size=11); ca = Alignment(vertical='center',wrap_text=True)
    bd = Border(left=Side(style='thin'),right=Side(style='thin'),top=Side(style='thin'),bottom=Side(style='thin'))
    alt_fill = PatternFill(start_color='F0F4FA', end_color='F0F4FA', fill_type='solid')

    # Парсим chart маркер
    chart_info = parse_excel_chart(content)
    # Убираем маркеры из контента
    clean_content = re.sub(r'\[EXCEL_CHART\].*?\[/EXCEL_CHART\]', '', content, flags=re.DOTALL)

    tl = [l.strip() for l in clean_content.split('\n') if l.strip() and '|' in l]
    if tl:
        rn = 0
        for line in tl:
            if all(c in '-| ' for c in line): continue
            rn += 1; cells = [c.strip() for c in line.split('|') if c.strip()]
            for ci, v in enumerate(cells, 1):
                cell = ws.cell(row=rn, column=ci, value=v); cell.border = bd
                if rn == 1: cell.font, cell.fill, cell.alignment = hf, hfl, ha
                else:
                    cell.font, cell.alignment = cf, ca
                    if rn % 2 == 0: cell.fill = alt_fill
                    # Пробуем преобразовать в число
                    try:
                        num = float(v.replace(' ', '').replace(',', '.'))
                        cell.value = num
                    except: pass
    # Автоширина
    for col in ws.columns:
        ml = max((len(str(c.value or '')) for c in col), default=10)
        ws.column_dimensions[col[0].column_letter].width = min(ml + 4, 50)

    # Добавляем диаграмму если есть числовые данные
    if chart_info and rn > 1:
        try:
            # Определяем числовые столбцы
            max_col = ws.max_column; max_row = ws.max_row
            if max_row > 1 and max_col > 1:
                ct = chart_info.get('type', 'bar')
                if ct == 'pie':
                    chart = PieChart()
                elif ct == 'line':
                    chart = LineChart()
                    chart.style = 10
                else:
                    chart = BarChart()
                    chart.style = 10
                chart.title = chart_info.get('title', 'Диаграмма')
                chart.width = 20; chart.height = 12

                data = Reference(ws, min_col=2, min_row=1, max_col=max_col, max_row=max_row)
                cats = Reference(ws, min_col=1, min_row=2, max_row=max_row)
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)

                ws_chart = wb.create_sheet("Диаграмма")
                ws_chart.add_chart(chart, "A1")
        except Exception as e:
            logger.error(f"Excel chart: {e}")

    buf = io.BytesIO(); wb.save(buf); buf.seek(0); return buf


# ─── Определение типа документа ───────────────────────────────

def detect_document_request(text):
    if '[PPTX_DOC]' in text:
        parts=text.split('[PPTX_DOC]',1);c=parts[1].strip() if len(parts)>1 else text;t="Презентация"
        for l in c.split('\n'):
            s=l.strip()
            if s.startswith('ЗАГОЛОВОК:'): t=s[10:].strip();break
            elif s.startswith('#'): t=s.lstrip('#').strip();break
            elif s and not s.startswith('['): t=s[:60];break
        return('pptx',t,c)
    elif '[WORD_DOC]' in text:
        parts=text.split('[WORD_DOC]',1);c=parts[1].strip() if len(parts)>1 else text;t="Материал"
        for l in c.split('\n'):
            if l.strip().startswith('#'): t=l.strip().lstrip('#').strip();break
            elif l.strip(): t=l.strip()[:60];break
        return('word',t,c)
    elif '[EXCEL_DOC]' in text:
        parts=text.split('[EXCEL_DOC]',1);c=parts[1].strip() if len(parts)>1 else text;t="Таблица"
        for l in c.split('\n'):
            if l.strip() and '|' not in l and not l.strip().startswith('['): t=l.strip()[:31];break
        return('excel',t,c)
    elif '[HTML_TEST]' in text:
        parts=text.split('[HTML_TEST]',1);c=parts[1].strip() if len(parts)>1 else text
        return('html','Интерактивный тест',c)
    return(None,'','')


# ─── Клавиатура (убрана кнопка Материал из меню) ──────────────

PERSISTENT_KEYBOARD = ReplyKeyboardMarkup(
    [[KeyboardButton("📋 Меню"), KeyboardButton("💰 Баланс")],
     [KeyboardButton("📽 Оформить презентацию"), KeyboardButton("📄 Оформить Word")],
     [KeyboardButton("📊 Оформить Excel"), KeyboardButton("🧪 Тест HTML")],
     [KeyboardButton("❓ Помощь")]],
    resize_keyboard=True, is_persistent=True
)


# ─── Генерация документа ─────────────────────────────────────

async def generate_document(update, context, uid, fmt, topic=None):
    if not check_generation_limit(uid):
        await safe_reply(update, GENERATION_LIMIT_MESSAGE); return
    h = get_history(uid); has = any(m['role']=='assistant' for m in h)
    prompts = {'pptx': PPTX_PROMPT, 'word': WORD_PROMPT, 'excel': EXCEL_PROMPT, 'html': HTML_TEST_PROMPT}
    labels = {'pptx': 'Создаю презентацию...', 'word': 'Создаю Word...', 'excel': 'Создаю Excel...', 'html': 'Создаю интерактивный тест...'}
    await safe_reply(update, labels.get(fmt, 'Создаю...'))
    base = f"На основе нашего диалога создай {fmt}." if has and not topic else f"Создай по теме: {topic or 'материал'}"
    answer = ask_claude(uid, base + prompts.get(fmt, ''))
    dt, dtitle, dcontent = detect_document_request(answer)
    if not dt:
        await safe_reply(update, "Не удалось. Обсудите тему подробнее и попробуйте снова."); return

    use_generation(uid)

    # Проверяем сохранённое фото для вставки в презентацию
    user_photo = context.user_data.get('last_photo_for_pptx')

    if dt=='pptx':
        photo_io = None
        if user_photo:
            photo_io = io.BytesIO(base64.b64decode(user_photo['base64']))
            context.user_data['last_photo_for_pptx'] = None
        buf=create_pptx_presentation(dtitle, dcontent, user_photo=photo_io)
        await update.message.reply_document(document=buf,filename=f"{dtitle[:30].replace(' ','_')}.pptx",caption=f"{dtitle}")
    elif dt=='word':
        buf=create_word_document(dtitle,dcontent)
        await update.message.reply_document(document=buf,filename=f"{dtitle[:30].replace(' ','_')}.docx",caption=f"{dtitle}")
    elif dt=='excel':
        buf=create_excel_document(dtitle,dcontent)
        await update.message.reply_document(document=buf,filename=f"{dtitle[:30].replace(' ','_')}.xlsx",caption=f"{dtitle}")
    elif dt=='html':
        # Отправляем HTML файл
        html_content = dcontent.strip()
        # Убираем markdown обёртки если есть
        if html_content.startswith('```'): html_content = re.sub(r'^```\w*\n?', '', html_content)
        if html_content.endswith('```'): html_content = html_content[:-3]
        html_buf = io.BytesIO(html_content.encode('utf-8'))
        html_buf.name = f"Тест_{dtitle[:20].replace(' ','_')}.html"
        await update.message.reply_document(document=html_buf, filename=html_buf.name, caption=f"Интерактивный тест - откройте в браузере!")


# ─── Обработчики ──────────────────────────────────────────────

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    save_user_info(uid, update.effective_user.first_name, update.effective_user.username)
    ql, gl = get_limit_status(uid); tariff = get_tariff(uid)
    kb = [
        [InlineKeyboardButton("План урока", callback_data="plan"), InlineKeyboardButton("Тест", callback_data="test")],
        [InlineKeyboardButton("Идеи", callback_data="ideas"), InlineKeyboardButton("ОГЭ/ЕГЭ", callback_data="ege")],
    ]
    w = (
        f"Добро пожаловать в {BOT_NAME}!\n\n"
        "ИИ-помощник для учителей любого предмета.\n\n"
        "Напишите что нужно:\n"
        "- План урока по биологии, 6 класс\n"
        "- Тест: дроби, 5 класс, 15 вопросов\n"
        "- Отправьте фото учебника\n"
        "- Запишите голосовое\n\n"
        "Когда готово - кнопки оформления внизу\n"
        "Новинка: кнопка Тест HTML - интерактивный тест для учеников!\n\n"
        f"Тариф: {tariff['name']}\n"
        f"Запросов: {ql} | Оформлений: {gl}"
    )
    await update.message.reply_text(w, reply_markup=InlineKeyboardMarkup(kb))
    await update.message.reply_text("Кнопки:", reply_markup=PERSISTENT_KEYBOARD)

async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id; ql, gl = get_limit_status(uid)
    await safe_reply(update,
        f"{BOT_NAME} v{BOT_VERSION}\n\n"
        "Как пользоваться:\n"
        "1. Напишите запрос\n"
        "2. Обсуждайте, уточняйте\n"
        "3. Нажмите кнопку оформления\n\n"
        "Форматы: презентация, Word, Excel, HTML тест\n\n"
        "Команды: /start /help /clear /balance\n"
        "/pptx тема - презентация\n"
        "/word тема - документ\n"
        "/excel тема - таблица\n\n"
        f"Запросов: {ql} | Оформлений: {gl}\n"
        f"Поддержка: {SUPPORT_BOT_URL}")

async def balance_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await safe_reply(update, get_balance_text(update.effective_user.id))

async def status_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await balance_command(update, context)

async def clear_history_cmd(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_conversations[update.effective_user.id] = []
    await safe_reply(update, "История очищена!")

async def pptx_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await generate_document(update, context, update.effective_user.id, 'pptx', ' '.join(context.args) if context.args else None)

async def word_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await generate_document(update, context, update.effective_user.id, 'word', ' '.join(context.args) if context.args else None)

async def excel_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await generate_document(update, context, update.effective_user.id, 'excel', ' '.join(context.args) if context.args else None)

BUTTON_PROMPTS = {
    "plan": "Тема, предмет, класс:\nПлан урока: Фотосинтез, биология, 6 класс",
    "test": "Тема, предмет, кол-во:\nТест: Дроби, 5 класс, 20 вопросов",
    "ege": "Предмет:\nРазбор задания ЕГЭ, русский, задание 8",
    "ideas": "Предмет:\nИдеи: Электричество, физика, 8 класс",
}

async def button_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query; await q.answer()
    if q.data in BUTTON_PROMPTS: await q.message.reply_text(BUTTON_PROMPTS[q.data])

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id
    if not check_query_limit(uid): await safe_reply(update, LIMIT_MESSAGE); return
    doc = update.message.document; caption = update.message.caption or ""; fn = doc.file_name or "file"
    ext = fn.rsplit('.',1)[-1].lower() if '.' in fn else ""
    if ext not in ['docx','doc','pdf','xlsx','xls','txt','csv']:
        await safe_reply(update, f".{ext} не поддерживается."); return
    await update.message.chat.send_action('typing'); await safe_reply(update, f"Читаю {fn}...")
    try:
        file = await doc.get_file(); fp = f"/tmp/doc_{uid}.{ext}"; await file.download_to_drive(fp); et = ""
        if ext=='docx':
            try:
                d=DocxDocument(fp);parts=[p.text for p in d.paragraphs if p.text.strip()]
                for t in d.tables:
                    for r in t.rows: rt=' | '.join(c.text.strip() for c in r.cells);parts.append(rt) if rt.strip() else None
                et='\n'.join(parts)
            except Exception as e: logger.error(f"docx:{e}")
        elif ext=='pdf':
            try:
                with open(fp,'rb') as f: reader=PyPDF2.PdfReader(f);et='\n'.join(pg.extract_text() or '' for pg in reader.pages)
            except Exception as e: logger.error(f"pdf:{e}")
        elif ext in ['xlsx','xls']:
            try:
                wb2=openpyxl.load_workbook(fp,read_only=True);rows=[]
                for sn in wb2.sheetnames:
                    ws2=wb2[sn];rows.append(f"--- {sn} ---")
                    for r in ws2.iter_rows(values_only=True): rt=' | '.join(str(c) if c else '' for c in r);rows.append(rt) if rt.strip() else None
                et='\n'.join(rows);wb2.close()
            except Exception as e: logger.error(f"xlsx:{e}")
        elif ext in ['txt','csv']:
            try:
                with open(fp,'r',encoding='utf-8') as f: et=f.read()
            except:
                with open(fp,'r',encoding='cp1251') as f: et=f.read()
        try: os.remove(fp)
        except: pass
        if not et or len(et.strip())<10: await safe_reply(update, "Не удалось извлечь текст."); return
        if len(et)>12000: et=et[:12000]+"\n\n[...обрезано...]"
        await safe_reply(update, f"Извлечено {len(et)} символов")
        if caption:
            use_query(uid);answer=ask_claude(uid,f"Документ {fn}:\n\n{et}\n\n{caption}")
            await safe_reply(update, clean_for_chat(answer))
        else:
            context.user_data['last_document']=et;context.user_data['last_document_name']=fn
            await safe_reply(update, "Файл прочитан! Что сделать?\n- Создай тест\n- План урока\n- Выбери ключевые темы")
    except Exception as e: logger.error(f"Doc:{e}");await safe_reply(update, "Ошибка при обработке файла.")

async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid=update.effective_user.id
    if not check_query_limit(uid): await safe_reply(update, LIMIT_MESSAGE);return
    await update.message.chat.send_action('typing');await safe_reply(update, "Распознаю...")
    try:
        v=update.message.voice or update.message.audio;f=await v.get_file()
        op=f"/tmp/voice_{uid}.ogg";wp=f"/tmp/voice_{uid}.wav";await f.download_to_drive(op)
        subprocess.run(['ffmpeg','-y','-i',op,'-ar','16000','-ac','1',wp],capture_output=True,timeout=30)
        rec=sr.Recognizer()
        with sr.AudioFile(wp) as src: audio=rec.record(src)
        text=rec.recognize_google(audio,language="ru-RU")
        try: os.remove(op);os.remove(wp)
        except: pass
        if not text: await safe_reply(update, "Не удалось распознать.");return
        await safe_reply(update, f"Распознано: {text}")
        use_query(uid);answer=ask_claude(uid,text);await safe_reply(update,clean_for_chat(answer))
    except sr.UnknownValueError: await safe_reply(update, "Не распознано.")
    except Exception as e: logger.error(f"Voice:{e}");await safe_reply(update, "Ошибка.")

async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid=update.effective_user.id
    if not check_query_limit(uid): await safe_reply(update, LIMIT_MESSAGE);return
    caption=update.message.caption or "";await update.message.chat.send_action('typing')
    photo=update.message.photo[-1];f=await photo.get_file();pb=await f.download_as_bytearray()
    idata={'base64':base64.b64encode(bytes(pb)).decode('utf-8'),'mime':'image/jpeg'}
    if not caption:
        await safe_reply(update, "Фото получено! Что сделать?\n- Создай тест по этой теме\n- Объясни что изображено\n- Составь план урока\n- Вставь в презентацию (потом нажмите Оформить презентацию)\n\nИли отправьте фото с подписью.")
        context.user_data['last_photo']=idata
        context.user_data['last_photo_for_pptx']=idata  # Сохраняем для вставки в PPTX
        return
    # Проверяем просьбу вставить в презентацию
    cl = caption.lower()
    if any(kw in cl for kw in ['вставь в презентац', 'добавь в презентац', 'в слайд']):
        context.user_data['last_photo_for_pptx'] = idata
        await safe_reply(update, "Фото сохранено! Теперь нажмите 'Оформить презентацию' и оно будет вставлено на титульный слайд.")
        return
    await safe_reply(update, "Анализирую...")
    use_query(uid);answer=ask_claude(uid,f"Изображение. {caption}\n\nПроанализируй и ответь.",image_data=idata)
    await safe_reply(update,clean_for_chat(answer))

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid=update.effective_user.id;text=update.message.text
    # Кнопки (без кнопки Материал)
    if text=="📋 Меню":
        kb=[[InlineKeyboardButton("План урока",callback_data="plan"),InlineKeyboardButton("Тест",callback_data="test")],
            [InlineKeyboardButton("Идеи",callback_data="ideas"),InlineKeyboardButton("ОГЭ/ЕГЭ",callback_data="ege")]]
        await update.message.reply_text("Выберите:",reply_markup=InlineKeyboardMarkup(kb));return
    if text=="💰 Баланс": await safe_reply(update, get_balance_text(uid));return
    if text=="📽 Оформить презентацию": await generate_document(update,context,uid,'pptx');return
    if text=="📄 Оформить Word": await generate_document(update,context,uid,'word');return
    if text=="📊 Оформить Excel": await generate_document(update,context,uid,'excel');return
    if text=="🧪 Тест HTML": await generate_document(update,context,uid,'html');return
    if text=="❓ Помощь": await help_command(update,context);return
    if not check_query_limit(uid): await safe_reply(update, LIMIT_MESSAGE);return
    await update.message.chat.send_action('typing')
    # Сохранённый документ
    sd=context.user_data.get('last_document')
    if sd:
        dn=context.user_data.get('last_document_name','документ');context.user_data['last_document']=None
        use_query(uid);answer=ask_claude(uid,f"Документ {dn}:\n\n{sd}\n\n{text}")
        await safe_reply(update,clean_for_chat(answer));return
    # Сохранённое фото
    sp=context.user_data.get('last_photo')
    if sp:
        context.user_data['last_photo']=None;await safe_reply(update, "Анализирую...")
        use_query(uid);answer=ask_claude(uid,f"Изображение. {text}",image_data=sp)
        await safe_reply(update,clean_for_chat(answer));return
    # Обычный запрос
    use_query(uid);answer=ask_claude(uid,text);await safe_reply(update,clean_for_chat(answer))

async def error_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    logger.error(f"Err:{context.error}")
    if update and update.message:
        try: await update.message.reply_text("Ошибка. Попробуйте /clear и повторите.")
        except: pass

def main():
    if not TELEGRAM_TOKEN or not ANTHROPIC_API_KEY: print("Нет токенов!");return
    print(f"{BOT_NAME} v{BOT_VERSION} запускается...")
    app=Application.builder().token(TELEGRAM_TOKEN).build()
    app.add_handler(CommandHandler("start",start));app.add_handler(CommandHandler("help",help_command))
    app.add_handler(CommandHandler("clear",clear_history_cmd));app.add_handler(CommandHandler("status",status_command))
    app.add_handler(CommandHandler("balance",balance_command))
    app.add_handler(CommandHandler("pptx",pptx_command));app.add_handler(CommandHandler("word",word_command));app.add_handler(CommandHandler("excel",excel_command))
    app.add_handler(CallbackQueryHandler(button_handler))
    app.add_handler(MessageHandler(filters.VOICE|filters.AUDIO,handle_voice))
    app.add_handler(MessageHandler(filters.Document.ALL,handle_document))
    app.add_handler(MessageHandler(filters.PHOTO,handle_photo))
    app.add_handler(MessageHandler(filters.TEXT&~filters.COMMAND,handle_message))
    app.add_error_handler(error_handler)
    print(f"{BOT_NAME} v{BOT_VERSION} запущен!")
    app.run_polling(allowed_updates=Update.ALL_TYPES)

if __name__=="__main__": main()
